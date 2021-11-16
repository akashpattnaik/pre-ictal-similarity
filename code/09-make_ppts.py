# %%
from pptx import Presentation
import json
from os.path import join as ospj
import glob
from PIL import Image
import pandas as pd
import os

def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size

    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width

    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)

    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side

# %%
# Get paths from config file and metadata
with open("config.json") as f:
    config = json.load(f)
repo_path = config['repositoryPath']
metadata_path = config['metadataPath']
palette = config['lightColors']

data_path = ospj(repo_path, 'data')
figure_path = ospj(repo_path, 'figures')

patient_cohort = pd.read_excel(ospj(data_path, "patient_cohort.xlsx"))

# %% Make subgraphs and expression ppt
prs = Presentation()
for index, row in patient_cohort.iterrows():
    pt = row["Patient"]
    pt = "HUP130"
    print(pt)

    pt_data_path = ospj(data_path, pt)
    pt_figure_path = ospj(figure_path, pt)
    if not os.path.exists(ospj(pt_figure_path, 'subgraphs_all.png')):
        print("No subgraphs figure, skipped subject")
        break

    slide = prs.slides.add_slide(prs.slide_layouts[8])
    _add_image(slide, 1, ospj(pt_figure_path, 'subgraphs_all.png'))
    title_placeholder = slide.shapes.title
    title_placeholder.text = pt

    for fname in glob.glob(ospj(pt_figure_path, 'subgraph_expression_sz_*_all.png')):
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        _add_image(slide, 1, fname)
        title_placeholder = slide.shapes.title
        title_placeholder.text = pt

    prs.save(ospj(figure_path, 'subgraphs_and_expression_all.pptx'))
    break
# %% Make sz similarity
prs = Presentation()
for index, row in patient_cohort.iterrows():
    pt = row["Patient"]
    pt = "HUP130"
    print(pt)

    pt_data_path = ospj(data_path, pt)
    pt_figure_path = ospj(figure_path, pt)

    slide = prs.slides.add_slide(prs.slide_layouts[8])
    if not os.path.exists(ospj(pt_figure_path, 'sz_dissim_mat_all.png')):
        print("No dissim mat fig, skipped subject")
        break

    _add_image(slide, 1, ospj(pt_figure_path, 'sz_dissim_mat_all.png'))
    title_placeholder = slide.shapes.title
    title_placeholder.text = pt

    prs.save(ospj(figure_path, 'sz_dissimilarities_all.pptx'))

    break
# %% Make subgraphs figure
n_components = 6
patient_cohort = pd.read_csv(ospj(data_path, "patient_cohort_with_soz_states.csv"))

prs = Presentation()
for index, row in patient_cohort.iterrows():
    pt = row["Patient"]
    pt = "HUP130"
    print(pt)
    if pt == "HUP073":
        continue
    pt_soz_state = row["SOZ Sensitive State (0-index)"]

    pt_data_path = ospj(data_path, pt)
    pt_figure_path = ospj(figure_path, pt)

    for i_comp in range(n_components):
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        _add_image(slide, 1, ospj(pt_figure_path, "soz_subgraph_{}_heatmap_all.png".format(i_comp)))
        title_placeholder = slide.shapes.title
        if i_comp == pt_soz_state:
            title_placeholder.text = "{}, soz state".format(pt)
        else:
            title_placeholder.text = pt
            
    for fname in glob.glob(ospj(pt_figure_path, 'soz_expression_sz_*_all.png')):
        slide = prs.slides.add_slide(prs.slide_layouts[8])
        _add_image(slide, 1, fname)
        title_placeholder = slide.shapes.title
        title_placeholder.text = pt

    break
prs.save(ospj(figure_path, 'subgraphs_and_soz_states_all.pptx'))
# %%
